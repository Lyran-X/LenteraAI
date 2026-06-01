from dataclasses import dataclass
from pathlib import Path


UNSUPPORTED_FILE_MESSAGE = "Unsupported file type. Please upload PDF, TXT, DOCX, or PPTX."


@dataclass(frozen=True)
class PageText:
    page_number: int
    text: str
    char_start: int
    char_end: int


@dataclass(frozen=True)
class ParsedDocument:
    text: str
    page_count: int
    pages: list[PageText]


class FileParser:
    supported_extensions = {".pdf", ".txt", ".docx", ".pptx"}
    supported_mime_types = {
        "application/pdf",
        "text/plain",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }

    def parse(self, file_path: str | Path, mime_type: str, filename: str) -> ParsedDocument:
        path = Path(file_path)
        extension = path.suffix.lower() or Path(filename).suffix.lower()
        normalized_mime_type = (mime_type or "").lower()

        if normalized_mime_type == "application/pdf" or extension == ".pdf":
            return self._parse_pdf(path)

        if normalized_mime_type == "text/plain" or extension == ".txt":
            return self._parse_txt(path)

        if (
            normalized_mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or extension == ".docx"
        ):
            return self._parse_docx(path)

        if (
            normalized_mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            or extension == ".pptx"
        ):
            return self._parse_pptx(path)

        raise ValueError(UNSUPPORTED_FILE_MESSAGE)

    def _parse_pdf(self, path: Path) -> ParsedDocument:
        try:
            import fitz
        except ImportError as exc:
            raise RuntimeError("PyMuPDF is not installed. Run pip install PyMuPDF.") from exc

        sections: list[tuple[int, str]] = []

        with fitz.open(path) as document:
            for index, page in enumerate(document, start=1):
                page_text = page.get_text("text").strip()
                if page_text:
                    sections.append((index, page_text))

        return self._build_parsed_document(sections)

    def _parse_txt(self, path: Path) -> ParsedDocument:
        raw_bytes = path.read_bytes()

        try:
            text = raw_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text = raw_bytes.decode("latin-1")

        text = text.strip()
        return self._build_parsed_document([(1, text)] if text else [])

    def _parse_docx(self, path: Path) -> ParsedDocument:
        try:
            from docx import Document
        except ImportError as exc:
            raise RuntimeError("python-docx is not installed. Run pip install python-docx.") from exc

        document = Document(path)
        sections: list[tuple[int, str]] = []
        section_number = 1

        for paragraph in document.paragraphs:
            paragraph_text = paragraph.text.strip()
            if paragraph_text:
                sections.append((section_number, paragraph_text))
                section_number += 1

        for table in document.tables:
            table_lines: list[str] = []
            for row in table.rows:
                cell_values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cell_values:
                    table_lines.append(" | ".join(cell_values))

            table_text = "\n".join(table_lines).strip()
            if table_text:
                sections.append((section_number, table_text))
                section_number += 1

        return self._build_parsed_document(sections)

    def _parse_pptx(self, path: Path) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("python-pptx is not installed. Run pip install python-pptx.") from exc

        presentation = Presentation(path)
        sections: list[tuple[int, str]] = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    slide_texts.append(text.strip())

            slide_text = "\n".join(slide_texts).strip()
            if slide_text:
                sections.append((slide_number, slide_text))

        return self._build_parsed_document(sections)

    def _build_parsed_document(self, sections: list[tuple[int, str]]) -> ParsedDocument:
        text_parts: list[str] = []
        pages: list[PageText] = []
        cursor = 0

        for page_number, section_text in sections:
            clean_text = section_text.strip()
            if not clean_text:
                continue

            if text_parts:
                text_parts.append("\n\n")
                cursor += 2

            start = cursor
            text_parts.append(clean_text)
            cursor += len(clean_text)
            pages.append(
                PageText(
                    page_number=page_number,
                    text=clean_text,
                    char_start=start,
                    char_end=cursor,
                )
            )

        full_text = "".join(text_parts).strip()
        return ParsedDocument(
            text=full_text,
            page_count=len(pages),
            pages=pages,
        )
